from pptx import Presentation

# Create a PowerPoint presentation object
prs = Presentation()

# Add a title slide
slide_layout = prs.slide_layouts[0]  # 0 corresponds to the title slide layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "My PowerPoint Presentation"
subtitle.text = "Created using Python and python-pptx"

# Add a content slide
slide_layout = prs.slide_layouts[1]  # 1 corresponds to the content slide layout
slide = prs.slides.add_slide(slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "Slide 2"
content = "This is the content of Slide 2."
p = body_shape.text_frame.add_paragraph()
p.text = content

# Save the PowerPoint presentation to a file
prs.save("my_presentation.pptx")
